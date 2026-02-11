import os

PROJECT_DIR = os.path.dirname(os.path.abspath(__file__))
os.environ["PATH"] = PROJECT_DIR + os.pathsep + os.environ.get("PATH", "")

import streamlit as st
import whisper
import google.generativeai as GenAI


from pptx import Presentation
from io import BytesIO
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Inches



from reportlab.lib.pagesizes import LETTER
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import Paragraph

import re


def load_whisper():
    return whisper.load_model("base")



# ======================================================================================================================
# PAGE CONFIG & UI HEADER
# ======================================================================================================================

st.set_page_config(page_title="Gen", page_icon="🪄")

st.markdown(
    """
    <h1 style="
        color: #FFFFFF;
        text-align: center;
        text-shadow: 2px 2px 10px rgba(0,0,0,0.7);
    ">
    🪄 𝑻𝒓𝒂𝒏𝒔𝒄𝒓𝒊𝒑𝒕𝒊𝒐𝒏 & 𝑺𝒍𝒊𝒅𝒆 𝑪𝒓𝒆𝒂𝒕𝒐𝒓
    </h1>
    """,
    unsafe_allow_html=True
)


# ======================================================================================================================
# BACKGROUND STYLING
# ======================================================================================================================

Url_Imagen = "https://i.pinimg.com/originals/cf/a2/39/cfa239195d194b724a9d38362859a1af.jpg"

st.markdown(
    f"""
    <style>
    .stApp {{
        background-image: url("{Url_Imagen}");
        background-size: cover;
        background-position: center;
        background-repeat: no-repeat;
        background-attachment: fixed;
    }}

    .main {{
        background-color: rgba(0, 0, 0, 0.45);
        padding: 20px;
        border-radius: 20px;
    }}
    </style>
    """,
    unsafe_allow_html=True
)


# ======================================================================================================================
# API CONFIGURATION
# ======================================================================================================================

GenAI.configure(api_key=st.secrets["API_KEY"])


# ======================================================================================================================
# POWERPOINT CREATION FUNCTION (MODIFICADA)
# ======================================================================================================================
# AHORA RECIBE EL ARGUMENTO: titulo_notas
def crear_pptx(texto_generado, titulo_notas="SPEAKER NOTES"):
    prs = Presentation("template.pptx")
    
    # Dimensiones estándar de una diapositiva (en puntos)
    SLIDE_WIDTH = prs.slide_width
    SLIDE_HEIGHT = prs.slide_height

    pattern = r"---\s*SLIDE\s*\d+\s*---\s*(.*?)\s*(?=(?:---\s*SLIDE\s*\d+\s*---)|\Z)"
    slides = re.findall(pattern, texto_generado, flags=re.S)

    if not slides:
        slides = [s for s in re.split(r"---\s*SLIDE", texto_generado) if s.strip()]

    for slide_text in slides:
        slide_text = slide_text.strip()
        if not slide_text: continue

        # Separación de contenido
        if "###NOTAS###" in slide_text:
            partes = slide_text.split("###NOTAS###")
            lineas = partes[0].strip().splitlines()
            title_text = lineas[0].strip() if lineas else "Sin Título"
            body_content = "\n".join(lineas[1:]).strip()
            notes_text = partes[1].strip()
        else:
            lineas = [l.strip() for l in slide_text.splitlines() if l.strip()]
            title_text = lineas[0] if lineas else "Sin Título"
            body_content = "\n".join(lineas[1:])
            notes_text = ""

        # Crear diapositiva limpia
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # --- CONFIGURACIÓN DEL TÍTULO ---
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.left = Pt(40)
            title_shape.top = Pt(20)
            title_shape.width = SLIDE_WIDTH - Pt(80) 
            title_shape.height = Pt(90) 
            
            title_shape.text = title_text
            tf_title = title_shape.text_frame
            tf_title.word_wrap = True
            tf_title.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 
            
            for paragraph in tf_title.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(28) 

        # --- CONFIGURACIÓN DEL CUERPO Y NOTAS ---
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            body_shape.left = Pt(40)
            body_shape.top = Pt(130) 
            body_shape.width = SLIDE_WIDTH - Pt(80)
            body_shape.height = SLIDE_HEIGHT - Pt(150) 

            tf = body_shape.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            # Párrafo principal
            p = tf.paragraphs[0]
            p.text = body_content
            p.font.size = Pt(16)
            
            if notes_text:
                tf.add_paragraph() # Espacio en blanco
                
                # Etiqueta Notas (USANDO LA VARIABLE DINÁMICA)
                p_label = tf.add_paragraph()
                p_label.text = f"➤ {titulo_notas}:" 
                p_label.font.bold = True
                p_label.font.size = Pt(12)
                p_label.font.color.rgb = RGBColor(100, 100, 100) 
                
                # Contenido Notas
                p_notes = tf.add_paragraph()
                p_notes.text = notes_text
                p_notes.font.italic = True
                p_notes.font.size = Pt(11)
                p_notes.font.color.rgb = RGBColor(120, 120, 120)

    pptx_io = BytesIO()
    prs.save(pptx_io)
    return pptx_io.getvalue()

# ======================================================================================================================
# PDF CREATION FUNCTION (MODIFICADA)
# ======================================================================================================================
# AHORA RECIBE EL ARGUMENTO: titulo_notas
def crear_pdf(texto_generado, titulo_notas="SPEAKER NOTES"):
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=LETTER)
    width, height = LETTER
    styles = getSampleStyleSheet()
    
    # Estilos
    body_style = ParagraphStyle('Body', fontSize=11, leading=14, spaceAfter=10)
    title_style = ParagraphStyle('Title', fontSize=20, leading=24, spaceAfter=20, fontName='Helvetica-Bold')
    label_style = ParagraphStyle('Label', fontSize=10, leading=12, fontName='Helvetica-BoldOblique', textColor='grey')
    note_content_style = ParagraphStyle('NoteBody', fontSize=9, leading=11, leftIndent=10, textColor='grey')

    slides = re.findall(r"---\s*SLIDE\s*\d+\s*---\s*(.*?)\s*(?=(?:---\s*SLIDE\s*\d+\s*---)|\Z)", texto_generado, flags=re.S)
    if not slides: slides = [s for s in re.split(r"---\s*SLIDE", texto_generado) if s.strip()]

    for slide_text in slides:
        if "###NOTAS###" in slide_text:
            partes = slide_text.split("###NOTAS###")
            content_lines = [l.strip() for l in partes[0].strip().splitlines() if l.strip()]
            title_text = content_lines[0]
            body_text = " ".join(content_lines[1:])
            notes_text = partes[1].strip()
        else:
            lines = [l.strip() for l in slide_text.splitlines() if l.strip()]
            title_text = lines[0]
            body_text = " ".join(lines[1:])
            notes_text = ""

        curr_y = height - 1*inch

        # Dibujar Título
        p_title = Paragraph(title_text, title_style)
        w, h = p_title.wrap(width - 2*inch, height)
        curr_y -= h
        p_title.drawOn(c, 1*inch, curr_y)

        # Dibujar Cuerpo
        curr_y -= 0.3*inch
        p_body = Paragraph(body_text, body_style)
        w, h = p_body.wrap(width - 2*inch, height)
        curr_y -= h
        p_body.drawOn(c, 1*inch, curr_y)

        # Dibujar Notas si existen
        if notes_text:
            curr_y -= 0.4*inch
            # Etiqueta Notas (USANDO LA VARIABLE DINÁMICA)
            p_label = Paragraph(f"{titulo_notas}:", label_style)
            w, h = p_label.wrap(width - 2*inch, height)
            curr_y -= h
            p_label.drawOn(c, 1*inch, curr_y)

            p_note = Paragraph(notes_text, note_content_style)
            w, h = p_note.wrap(width - 2.2*inch, height)
            curr_y -= h
            p_note.drawOn(c, 1.2*inch, curr_y)

        c.showPage()

    c.save()
    buffer.seek(0)
    return buffer.getvalue()

# ======================================================================================================================
# AUDIO UPLOAD & TRANSCRIPTION
# ======================================================================================================================
audio_Recorded = st.audio_input("𝑹𝒆𝒄𝒐𝒓𝒅 𝒚𝒐𝒖𝒓 𝒂𝒖𝒅𝒊𝒐")
audio_Fill = st.file_uploader(
    "𝑼𝒑𝒍𝒐𝒂𝒅 𝒚𝒐𝒖𝒓 𝒂𝒖𝒅𝒊𝒐",
    type=["mp3", "mp4", "opus", "wav", "m4a"]
)
Audio_fill = audio_Fill or audio_Recorded


if Audio_fill is not None:
    MAX_FILE_SIZE = 30 * 1024 * 1024
    if Audio_fill.size > MAX_FILE_SIZE:
        st.error("The audio is too long or too short. Please upload a file shorter than 30 minutes. (MAX 30MB)")
        st.stop()
    with st.expander("Show audio"):
     st.audio(Audio_fill)

    with open("temp_audio.wav", "wb") as f:
        f.write(Audio_fill.getbuffer())

    with st.spinner("Whisper is processing your audio"):
        modelo_whisper = load_whisper()
        resultado = modelo_whisper.transcribe("temp_audio.wav")
        
        #DETECCIÓN DE IDIOMA 
        detected_language_code = resultado.get("language", "en")


    with st.expander("Show transcription"):
        st.write(f"Detected Language: {detected_language_code}")
        st.write(resultado["text"])


# ======================================================================================================================
# GENERATIVE SLIDES
# ======================================================================================================================

if Audio_fill is not None and st.button("✨ Generative Slides"):

    # Diccionario simple con los idiomas más comunes, por defecto usa Inglés
    translations_map = {
        "es": "NOTAS DE ORADOR",
        "en": "SPEAKER NOTES",
        "fr": "NOTES DE L'ORATEUR",
        "de": "SPRECHERNOTIZEN",
        "it": "NOTE DEL RELATORE",
        "pt": "NOTAS DO ORADOR",
        "ru": "ЗАМЕТКИ ДОКЛАДЧИКА",
        "zh": "演讲者备注",
        "ja": "スピーカーノート"
    }
    
    # Selecciona la traducción basada en el código de Whisper 
    titulo_notas_final = translations_map.get(detected_language_code, "SPEAKER NOTES")


    with st.spinner("Gemini is creating your slides..."):
        modelo_gemini = GenAI.GenerativeModel('models/gemini-2.5-flash')

        instruction = f"""
        CORE MISSION: Analyze the provided transcription and generate a professional presentation based ONLY on its content.

        [STRICT LANGUAGE RULE]

        Identify the language of the input: {resultado['text']}.

        ALL generated content MUST be in that exact language.

        DO NOT translate. If the input is Spanish, the output is 100% Spanish.

        STRATEGIC AGENT ROLE:

        Act as a Strategic Consultant. Transform the transcript into a high-level corporate/academic narrative.

        NO META-COMMENTARY: Do not use phrases like "The transcript says" or "This slide covers." State information as objective, established facts.

        DEVELOPMENT: Expand brief mentions into sophisticated, professional concepts.

        OUTPUT FORMAT (MANDATORY STRUCTURE): Generate at least 5 slides. Use the following structure for each one, but DO NOT include labels like "Title:", "Text:", or "Notes:".

        --- SLIDE [N] ---

        [Insert Professional Title Here]
        [Insert here a single paragraph of 4 to 6 lines. Use fluid and professional prose. STRICTLY PROHIBITED: Bullet points, lists, or internal labels.]

        ###NOTAS###
        [Insert here a professional script for the speaker in the SAME language as the transcript. Explain the "why" behind the concept and its strategic connection to the next slide.]

        STRICT CONSTRAINTS:

        NO LABELS: Do not write the words "Title", "Text", "Slide", or "notes_slide" inside the content. Use only the Markdown Header (▶) for the title.

        NO LISTS: Use only continuous paragraphs.

        SILENT EXECUTION: Do not include greetings, introductions, or conclusions (e.g., "Here are your slides"). Return ONLY the structured slide content.

        MINIMUM VOLUME: If the transcript is short, use logical professional deduction to reach exactly 5 slides.

        INPUT DATA: {resultado['text']}
                """

        answer = modelo_gemini.generate_content(instruction)


    with st.expander("Show Content"):
        st.write(answer.text)

    #PASAMOS EL TÍTULO TRADUCIDO A LAS FUNCIONES
    pptx_data = crear_pptx(answer.text, titulo_notas_final)
    pdf_data = crear_pdf(answer.text, titulo_notas_final)

    st.download_button(
        label="🚀 DOWNLOAD YOUR POWERPOINT",
        data=pptx_data,
        file_name="Presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True
    )


    st.download_button(
    label="📄 DOWNLOAD PDF",
    data=pdf_data,
    file_name="Presentation.pdf",
    mime="application/pdf",
    use_container_width=True
)


    st.balloons()

    if os.path.exists("temp_audio.wav"):
        os.remove("temp_audio.wav")
